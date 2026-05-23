#!/usr/bin/env python3
"""Generate slide deck: Three-Layer Agent Architecture"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# Colors
BG_DARK = RGBColor(0x0A, 0x0A, 0x0F)
SURFACE = RGBColor(0x12, 0x12, 0x1A)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
TEXT_WHITE = RGBColor(0xE4, 0xE4, 0xE7)
TEXT_MUTED = RGBColor(0x71, 0x71, 0x7A)
BORDER = RGBColor(0x1E, 0x1E, 0x2E)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_card(slide, left, top, width, height, title, items, accent_color=ACCENT_PURPLE):
    """Add a card/box with title and bullet items."""
    # Background rectangle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    # Accent bar at top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Title
    add_text_box(slide, left + 0.3, top + 0.2, width - 0.6, 0.5, title, font_size=16, color=accent_color, bold=True)

    # Items
    y = top + 0.7
    for item in items:
        add_text_box(slide, left + 0.3, y, width - 0.6, 0.35, f"→ {item}", font_size=12, color=TEXT_WHITE)
        y += 0.35

    return shape


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
set_slide_bg(slide)

add_text_box(slide, 1, 1.5, 11, 1, "Craftura.ai", font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)
add_text_box(slide, 1, 2.5, 11, 1.5, "Three-Layer Agent Architecture", font_size=44, color=TEXT_WHITE, bold=True)
add_text_box(slide, 1, 4.2, 11, 0.8, "Technical Reference with Code-Level Evidence and Industry Citations", font_size=20, color=ACCENT_PURPLE)

# Three colored dots representing layers
for i, (x, color) in enumerate([(5.5, ACCENT_PURPLE), (6.5, ACCENT_CYAN), (7.5, ACCENT_AMBER)]):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(5.8), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

add_text_box(slide, 1, 6.2, 11, 0.5, "Runtime Orchestration  •  LLM Reasoning  •  Memory Storage", font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: The Core Argument
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "THE CORE ARGUMENT", font_size=32, color=TEXT_WHITE, bold=True)

# Main statement box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(11.3), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = SURFACE
shape.line.color.rgb = ACCENT_PURPLE
shape.line.width = Pt(2)

add_text_box(slide, 1.5, 1.8, 10.3, 1.8,
    "Supabase is NOT Jade's brain.\n\n"
    "Supabase is a PostgreSQL database that stores structured data with ACID guarantees.\n"
    "It does not reason, decide, classify, generate, or orchestrate.\n"
    "It stores and retrieves rows. That is the complete extent of its functionality.",
    font_size=18, color=TEXT_WHITE)

# Three key points
add_text_box(slide, 1, 4.5, 3.5, 0.4, "THE CLAIM", font_size=12, color=RED, bold=True)
add_text_box(slide, 1, 5.0, 3.5, 1.2, '"Supabase is Jade\'s brain"\n\nTechnically incorrect: conflates storage with cognition', font_size=14, color=TEXT_WHITE)

add_text_box(slide, 5, 4.5, 3.5, 0.4, "THE REALITY", font_size=12, color=GREEN, bold=True)
add_text_box(slide, 5, 5.0, 3.5, 1.2, '"Supabase stores Jade\'s memory"\n\nTechnically correct: database provides persistent cross-session storage', font_size=14, color=TEXT_WHITE)

add_text_box(slide, 9, 4.5, 3.5, 0.4, "THE IMPLICATION", font_size=12, color=ACCENT_AMBER, bold=True)
add_text_box(slide, 9, 5.0, 3.5, 1.2, 'All three layers must be deployed and connected.\n\nA database alone does nothing without runtime + LLM.', font_size=14, color=TEXT_WHITE)


# ============================================================
# SLIDE 3: The Three Layers Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "THE THREE LAYERS", font_size=32, color=TEXT_WHITE, bold=True)

# Layer 1 Card
add_card(slide, 0.5, 1.5, 3.8, 4.5, "LAYER 1: AGENT RUNTIME", [
    "Python code that runs agents",
    "Routes work between specialists",
    "Manages pipeline state (StudioState)",
    "Calls tools, APIs, databases",
    "Enforces guardrails & cost controls",
    "Your code: CrafturaFlow, CrafturaCrew",
], ACCENT_PURPLE)

# Layer 2 Card
add_card(slide, 4.8, 1.5, 3.8, 4.5, "LAYER 2: LLM PROCESSING", [
    "Stateless reasoning via API",
    "Analyzes, generates, decides",
    "Each call is independent",
    "No memory between calls",
    "Rented compute intelligence",
    "Your stack: GPT-4, Claude, Gemini, llama.cpp",
], ACCENT_CYAN)

# Layer 3 Card
add_card(slide, 9.1, 1.5, 3.8, 4.5, "LAYER 3: DATA & MEMORY", [
    "Persistent storage across sessions",
    "Structured data with ACID guarantees",
    "Identity, projects, conversations",
    "Audit trails & references",
    "Row-level security for access control",
    "Your schema: brain_profile, brain_projects, etc.",
], ACCENT_AMBER)


# ============================================================
# SLIDE 4: Code Evidence — Runtime Layer
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "CODE EVIDENCE: LAYER 1 — RUNTIME", font_size=28, color=ACCENT_PURPLE, bold=True)

# File reference
add_text_box(slide, 0.5, 1.2, 12, 0.4, "craftura-agents/src/craftura_agents/main.py", font_size=12, color=TEXT_MUTED)

# Code block
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(7.5), Inches(4.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x14)
shape.line.color.rgb = BORDER

code_text = (
    "class CrafturaFlow(Flow[StudioState]):\n"
    "    \"\"\"Jade Executive Assistant - Studio Orchestrator\"\"\"\n\n"
    "    @start()\n"
    "    def receive_lead(self):\n"
    "        self.state.current_stage = \"intake\"\n\n"
    "    @listen(\"classify_and_intake\")\n"
    "    def route_by_project_type(self):\n"
    "        if self.state.project_type in (\"existing_site\", \"redesign\"):\n"
    "            self.state.current_stage = \"audit\"\n\n"
    "    @listen(\"franklin_build\")\n"
    "    def run_qa(self):\n"
    "        result = self._get_crew().qa_crew().kickoff(inputs={...})\n"
)

add_text_box(slide, 0.8, 2.0, 6.9, 4.0, code_text, font_size=11, color=TEXT_WHITE, font_name="Consolas")

# Analysis box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.8), Inches(4.3), Inches(4.5))
shape.fill.solid()
shape.fill.fore_color.rgb = SURFACE
shape.line.color.rgb = ACCENT_PURPLE

add_text_box(slide, 8.8, 2.0, 3.7, 0.4, "WHAT THIS CODE DOES:", font_size=12, color=ACCENT_PURPLE, bold=True)
analysis = (
    "✓ Manages StudioState (30+ fields)\n\n"
    "✓ Routes leads through 10-stage pipeline\n\n"
    "✓ Enforces stage dependencies\n\n"
    "✓ Triggers escalation when needed\n\n"
    "✓ Calls crew methods for each stage\n\n"
    "✗ Does NOT perform reasoning\n"
    "✗ Does NOT store data\n"
    "✗ Is pure orchestration logic"
)
add_text_box(slide, 8.8, 2.5, 3.7, 3.5, analysis, font_size=11, color=TEXT_WHITE)


# ============================================================
# SLIDE 5: Code Evidence — LLM Layer
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "CODE EVIDENCE: LAYER 2 — LLM PROCESSING", font_size=28, color=ACCENT_CYAN, bold=True)

add_text_box(slide, 0.5, 1.2, 12, 0.4, "craftura-agents/src/craftura_agents/crews/craftura_crew/craftura_crew.py", font_size=12, color=TEXT_MUTED)

shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(7.5), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x14)
shape.line.color.rgb = BORDER

code_text2 = (
    "@agent\n"
    "def franklin(self) -> Agent:\n"
    "    return Agent(\n"
    "        config=self.agents_config[\"franklin\"],\n"
    "        verbose=True,\n"
    "        tools=[get_design_kit_reader_tool(),\n"
    "               get_pixabay_tool(),\n"
    "               get_flux_image_tool()],\n"
    "    )\n\n"
    "@crew\n"
    "def franklin_crew(self) -> Crew:\n"
    "    return Crew(\n"
    "        agents=[self.franklin()],\n"
    "        tasks=[self.franklin_implementation_task()],\n"
    "        process=Process.sequential,\n"
    "    )\n\n"
    "# CrewAI documentation:\n"
    "# \"Memory is a separate component configured\n"
    "#  independently from agents and crews.\"\n"
)

add_text_box(slide, 0.8, 2.0, 6.9, 3.0, code_text2, font_size=11, color=TEXT_WHITE, font_name="Consolas")

# Analysis box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.8), Inches(4.3), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = SURFACE
shape.line.color.rgb = ACCENT_CYAN

add_text_box(slide, 8.8, 2.0, 3.7, 0.4, "WHAT THIS CODE DOES:", font_size=12, color=ACCENT_CYAN, bold=True)
analysis2 = (
    "✓ Defines agents with system prompts\n\n"
    "✓ Attaches tools for capability\n\n"
    "✓ Sends prompts to LLM APIs\n\n"
    "✓ Receives reasoning output\n\n"
    "✓ Memory is SEPARATE (CrewAI docs)\n\n"
    "✗ Agents ≠ memory\n"
    "✗ Agents ≠ storage\n"
    "✗ LLM is stateless by design"
)
add_text_box(slide, 8.8, 2.5, 3.7, 2.7, analysis2, font_size=11, color=TEXT_WHITE)

# CrewAI citation
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = SURFACE
shape.line.color.rgb = ACCENT_CYAN

add_text_box(slide, 0.8, 5.8, 11.7, 0.8,
    'CrewAI Documentation: "Memory is a separate component configured independently from agents and crews."\n'
    'Source: https://docs.crewai.com/en/concepts/memory',
    font_size=12, color=TEXT_MUTED)


# ============================================================
# SLIDE 6: Code Evidence — Memory Layer
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "CODE EVIDENCE: LAYER 3 — MEMORY (SUPABASE)", font_size=28, color=ACCENT_AMBER, bold=True)

add_text_box(slide, 0.5, 1.2, 12, 0.4, "jade-assistant/.../brain/jade-brain-v1.sql", font_size=12, color=TEXT_MUTED)

shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(7.5), Inches(4.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x14)
shape.line.color.rgb = BORDER

code_text3 = (
    "create table if not exists public.brain_profile (\n"
    "    id uuid primary key default gen_random_uuid(),\n"
    "    section text not null unique,\n"
    "    content jsonb not default '{}'::jsonb,\n"
    ");\n\n"
    "create or replace function public.load_brain()\n"
    "returns jsonb language sql stable\n"
    "as $$\n"
    "  -- Returns: identity, operating_style,\n"
    "  -- recent_sessions, open_action_items,\n"
    "  -- recent_references, active_projects\n"
    "  select jsonb_build_object(\n"
    "    'identity', (select content from brain_profile\n"
    "                where section = 'identity'),\n"
    "    ...\n"
    "  );\n"
    "$$;\n\n"
    "-- Tables: brain_profile, brain_projects,\n"
    "--         brain_conversations, brain_references,\n"
    "--         brain_action_log\n"
)

add_text_box(slide, 0.8, 2.0, 6.9, 3.8, code_text3, font_size=11, color=TEXT_WHITE, font_name="Consolas")

# Analysis box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.8), Inches(4.3), Inches(4.2))
shape.fill.solid()
shape.fill.fore_color.rgb = SURFACE
shape.line.color.rgb = ACCENT_AMBER

add_text_box(slide, 8.8, 2.0, 3.7, 0.4, "WHAT THIS CODE IS:", font_size=12, color=ACCENT_AMBER, bold=True)
analysis3 = (
    "✓ PostgreSQL tables with JSONB\n\n"
    "✓ Standard SQL queries and indexes\n\n"
    "✓ Row-Level Security policies\n\n"
    "✓ Audit logging triggers\n\n"
    "✓ Data retrieval function\n\n"
    "✗ Zero AI capability\n"
    "✗ Zero reasoning\n"
    "✗ Zero decision-making\n"
    "✗ Standard database schema"
)
add_text_box(slide, 8.8, 2.5, 3.7, 3.4, analysis3, font_size=11, color=TEXT_WHITE)

# Supabase citation
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8))
shape.fill.solid()
shape.fill.fore_color.rgb = SURFACE
shape.line.color.rgb = ACCENT_AMBER

add_text_box(slide, 0.8, 6.45, 11.7, 0.5,
    'Supabase: "Most agent stacks require a vector database, an auth provider, a file store, an API layer, and a separate Postgres instance."\n'
    'Source: https://supabase.com/solutions/agents — Positioned as data infrastructure, not intelligence.',
    font_size=12, color=TEXT_MUTED)


# ============================================================
# SLIDE 7: How the Layers Connect
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "HOW THE LAYERS CONNECT", font_size=32, color=TEXT_WHITE, bold=True)

# Flow diagram using shapes and arrows
steps = [
    ("TRIGGER", "Lead / Command\nSchedule / Webhook", TEXT_MUTED, BORDER),
    ("LAYER 1", "Agent Runtime\nCrafturaFlow", ACCENT_PURPLE, ACCENT_PURPLE),
    ("↔", "", TEXT_WHITE, None),
    ("LAYER 2", "LLM API\nGPT-4 / Claude", ACCENT_CYAN, ACCENT_CYAN),
    ("↔", "", TEXT_WHITE, None),
    ("LAYER 3", "Supabase\nPostgres DB", ACCENT_AMBER, ACCENT_AMBER),
]

x = 0.5
for step in steps:
    if step[0] == "↔":
        add_text_box(slide, x, 2.8, 0.5, 0.5, "↔", font_size=24, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
        x += 0.5
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5), Inches(2.2), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = step[3]
        shape.line.width = Pt(2)

        add_text_box(slide, x + 0.15, 2.55, 1.9, 0.3, step[0], font_size=10, color=step[2], bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.15, 2.9, 1.9, 0.7, step[1], font_size=13, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
        x += 2.7

# Detailed flow steps below
flow_steps = [
    ("1.", "Runtime receives trigger (lead, command, schedule)", ACCENT_PURPLE),
    ("2.", "Runtime calls load_brain() on Supabase — SQL query returns JSONB", ACCENT_AMBER),
    ("3.", "Runtime injects loaded data into prompt + current task description", TEXT_WHITE),
    ("4.", "Runtime sends enriched prompt to LLM API", ACCENT_CYAN),
    ("5.", "LLM performs reasoning over the context and returns structured output", ACCENT_CYAN),
    ("6.", "Runtime receives LLM output, executes tools, takes actions", ACCENT_PURPLE),
    ("7.", "Runtime saves important results back to Supabase via write functions", ACCENT_AMBER),
]

y = 4.2
for num, desc, color in flow_steps:
    add_text_box(slide, 1, y, 0.5, 0.35, num, font_size=14, color=color, bold=True)
    add_text_box(slide, 1.5, y, 10, 0.35, desc, font_size=13, color=TEXT_WHITE)
    y += 0.38


# ============================================================
# SLIDE 8: Industry Citations — Framework Documentation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "INDUSTRY CITATIONS: FRAMEWORK DOCUMENTATION", font_size=28, color=TEXT_WHITE, bold=True)

citations = [
    ("CrewAI Memory Documentation",
     "Memory is a separate component configured independently from agents and crews.",
     "https://docs.crewai.com/en/concepts/memory"),
    ("LangChain Memory Overview",
     "Memory storage uses external stores. InMemoryStore saves to dictionary; use DB-backed store in production.",
     "https://docs.langchain.com/oss/python/concepts/memory"),
    ("MongoDB + LangGraph Integration",
     "Agent memory is a computational exocortex — integrates LLM context with persistent memory management system.",
     "https://www.mongodb.com/company/blog/product-release-announcements/powering-long-term-memory-for-agents-langgraph"),
    ("Supabase Agent Solutions",
     "Most agent stacks require vector DB, auth, file store, API layer, separate Postgres. Supabase replaces all of them.",
     "https://supabase.com/solutions/agents"),
]

y = 1.3
for title, quote, url in citations:
    # Citation box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = BORDER

    add_text_box(slide, 0.8, y + 0.05, 11.7, 0.3, title, font_size=14, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, 0.8, y + 0.4, 11.7, 0.4, f'"{quote}"', font_size=12, color=TEXT_WHITE)
    add_text_box(slide, 0.8, y + 0.85, 11.7, 0.3, url, font_size=10, color=TEXT_MUTED)
    y += 1.4


# ============================================================
# SLIDE 9: Industry Citations — Architecture Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "INDUSTRY CITATIONS: ARCHITECTURE ANALYSIS", font_size=28, color=TEXT_WHITE, bold=True)

citations2 = [
    ("Thinking.inc — Agentic AI Architecture (2025)",
     "The architecture required for agent systems — tool registries, orchestration layers, memory systems — goes well beyond what a generative AI application needs.",
     "https://thinking.inc/en/pillar-pages/agentic-ai-architecture"),
    ("Kore.ai — Multi-Agent Orchestration (2025)",
     "Multi-agent orchestration is the coordinated management of multiple AI agents so they work together as a unified, goal-driven system.",
     "https://www.kore.ai/blog/what-is-multi-agent-orchestration"),
    ("Dev.to — LangGraph Memory Walkthrough",
     "The LLM is stateless by design. Memory is always an infrastructure concern. The model only knows what is in the context window at inference time.",
     "https://dev.to/sreeni5018/five-agent-memory-types-in-langgraph-a-deep-code-walkthrough-part-2-17kb"),
    ("Menlo Ventures GenAI Stack (via Medium)",
     "The orchestration layer connects foundation models, data systems, and observability tools to create robust AI applications.",
     "https://medium.com/@akankshasinha247/agent-orchestration-when-to-use-langchain-langgraph-autogen-or-build-an-agentic-rag-system-cc298f785ea4"),
]

y = 1.3
for title, quote, url in citations2:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = BORDER

    add_text_box(slide, 0.8, y + 0.05, 11.7, 0.3, title, font_size=14, color=ACCENT_CYAN, bold=True)
    add_text_box(slide, 0.8, y + 0.4, 11.7, 0.4, f'"{quote}"', font_size=12, color=TEXT_WHITE)
    add_text_box(slide, 0.8, y + 0.85, 11.7, 0.3, url, font_size=10, color=TEXT_MUTED)
    y += 1.4


# ============================================================
# SLIDE 10: What "Supabase Is the Brain" Gets Wrong
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'WHAT "SUPABASE IS THE BRAIN" GETS WRONG', font_size=28, color=RED, bold=True)

# Comparison table
headers = ["CLAIM", "WHY IT'S TECHNICALLY INCORRECT", "WHAT'S ACTUALLY TRUE"]
col_widths = [3.5, 4.5, 4.3]
col_starts = [0.5, 4.2, 8.9]

# Header row
y = 1.3
for i, (header, width, start) in enumerate(zip(headers, col_widths, col_starts)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(start), Inches(y), Inches(width), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    shape.line.color.rgb = BORDER
    add_text_box(slide, start + 0.2, y + 0.05, width - 0.4, 0.4, header, font_size=11, color=TEXT_MUTED, bold=True)

rows = [
    ('"Supabase is Jade\'s brain"', 'A database does not think, reason, or make decisions', 'Supabase stores data that the runtime loads and sends to the LLM'),
    ('"Jade lives in Supabase"', 'Code runs as a process, not inside a database', "Jade's runtime code loads memory from Supabase at session start"),
    ('"The brain loads context"', 'Databases don\'t load context into anything', 'The runtime calls load_brain() and injects results into prompts'),
    ('"Brain functions think"', 'SQL functions return data; they do not reason', 'load_brain() returns JSONB. The LLM reasons over that data.'),
]

y = 1.8
for claim, wrong, true in rows:
    for i, (text, width, start) in enumerate(zip([claim, wrong, true], col_widths, col_starts)):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(start), Inches(y), Inches(width), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = BORDER
        color = RED if i == 0 else (TEXT_MUTED if i == 1 else GREEN)
        add_text_box(slide, start + 0.2, y + 0.1, width - 0.4, 0.7, text, font_size=11, color=color)
    y += 0.95


# ============================================================
# SLIDE 11: The Library Analogy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "THE CORRECT ANALOGY", font_size=32, color=TEXT_WHITE, bold=True)

# Three analogy boxes
analogies = [
    ("SUPABASE = LIBRARY", ACCENT_AMBER, [
        "Stores vast amounts of information",
        "Organized for efficient retrieval",
        "Does NOT answer questions itself",
        "A library ≠ a librarian's brain",
    ]),
    ("LLM = LIBRARIAN'S BRAIN", ACCENT_CYAN, [
        "Processes and reasons over information",
        "Finds relevant data from the library",
        "Synthesizes answers from retrieved facts",
        "Stateless — forgets between questions",
    ]),
    ("RUNTIME = LIBRARIAN'S BODY", ACCENT_PURPLE, [
        "Goes to the library and retrieves books",
        "Delivers answers to patrons",
        "Takes action based on findings",
        "Coordinates the whole process",
    ]),
]

x = 0.5
for title, color, items in analogies:
    add_card(slide, x, 1.8, 3.8, 4, title, items, color)
    x += 4.2


# ============================================================
# SLIDE 12: Deployment Reality
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "DEPLOYMENT REALITY", font_size=32, color=TEXT_WHITE, bold=True)

# Three scenarios
scenarios = [
    ("ONLY SUPABASE (Layer 3)", RED, [
        "Database with tables and functions exists",
        "No code is querying them",
        "No LLM is reasoning over the data",
        "Result: Dead data. Nothing happens.",
    ]),
    ("ONLY RUNTIME (Layer 1)", ACCENT_PURPLE, [
        "Process starts, calls load_brain()",
        "Gets connection error — no Supabase",
        "Falls back to local prompts only",
        "Result: Stateless agent. No persistent memory.",
    ]),
    ("ONLY LLM (Layer 2)", ACCENT_CYAN, [
        "Can have conversations with AI model",
        "Stateless — forgets between calls",
        "No persistent project tracking",
        "Result: Isolated chat completions. No agent system.",
    ]),
]

x = 0.5
for title, color, items in scenarios:
    add_card(slide, x, 1.5, 3.8, 4, title, items, color)
    x += 4.2

# Bottom statement
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = SURFACE
shape.line.color.rgb = GREEN
shape.line.width = Pt(2)

add_text_box(slide, 0.8, 6.0, 11.7, 0.8,
    "ALL THREE LAYERS MUST BE DEPLOYED AND CONNECTED.\n"
    "The runtime is the glue that makes the system functional.",
    font_size=16, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13: Full Citation List
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "FULL CITATION LIST", font_size=28, color=TEXT_WHITE, bold=True)

all_citations = [
    ("1.", "CrewAI Memory Documentation", "https://docs.crewai.com/en/concepts/memory"),
    ("2.", "LangChain Memory Overview", "https://docs.langchain.com/oss/python/concepts/memory"),
    ("3.", "MongoDB + LangGraph Integration", "https://www.mongodb.com/company/blog/product-release-announcements/powering-long-term-memory-for-agents-langgraph"),
    ("4.", "Thinking.inc — Agentic AI Architecture", "https://thinking.inc/en/pillar-pages/agentic-ai-architecture"),
    ("5.", "Kore.ai — Multi-Agent Orchestration", "https://www.kore.ai/blog/what-is-multi-agent-orchestration"),
    ("6.", "Supabase — Agent Solutions", "https://supabase.com/solutions/agents"),
    ("7.", "Medium — Cognitive Orchestration Layer", "https://medium.com/@raktims2210/cognitive-orchestration-layer-the-next-enterprise-ai-architecture-that-lets-hundreds-of-agents-35dd427811f3"),
    ("8.", "Menlo Ventures GenAI Stack (via Medium)", "https://medium.com/@akankshasinha247/agent-orchestration-when-to-use-langchain-langgraph-autogen-or-build-an-agentic-rag-system-cc298f785ea4"),
    ("9.", "LinkedIn — AI's Memory Bottleneck (Fivetran report)", "https://www.linkedin.com/posts/tobiemorganhitchcock_the-ai-race-is-now-about-databases-not-activity-7341483033788633088-jfaF"),
    ("10.", "Dev.to — LangGraph Memory Walkthrough", "https://dev.to/sreeni5018/five-agent-memory-types-in-langgraph-a-deep-code-walkthrough-part-2-17kb"),
]

y = 1.3
for num, title, url in all_citations:
    add_text_box(slide, 0.5, y, 0.4, 0.35, num, font_size=12, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, 1, y, 5, 0.35, title, font_size=12, color=TEXT_WHITE)
    add_text_box(slide, 6.5, y, 6.3, 0.35, url, font_size=11, color=TEXT_MUTED)
    y += 0.4


# ============================================================
# SLIDE 14: Conclusion
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "CONCLUSION", font_size=32, color=TEXT_WHITE, bold=True)

# Main conclusion box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(11.3), Inches(4.5))
shape.fill.solid()
shape.fill.fore_color.rgb = SURFACE
shape.line.color.rgb = ACCENT_PURPLE
shape.line.width = Pt(2)

conclusion = (
    "The three-layer architecture — Runtime (orchestration), LLM (reasoning), Memory (storage) — is:\n\n"
    "1.  How the code actually works — verified by reading main.py, craftura_crew.py, and jade-brain-v1.sql\n\n"
    "2.  How CrewAI is designed — memory is explicitly a separate component from agents\n\n"
    "3.  How LangChain/LangGraph is designed — memory stores are external to graph orchestration\n\n"
    "4.  How Supabase positions itself — as data infrastructure for agents, not as intelligence\n\n"
    "5.  How every production agent system works — orchestration, reasoning, and memory are distinct layers\n\n"
    "The best option — the most cited, most technically clean option — is to recognize all three layers as "
    "distinct, deploy all three layers, and connect them properly through the runtime orchestration layer."
)

add_text_box(slide, 1.5, 1.8, 10.3, 4.0, conclusion, font_size=15, color=TEXT_WHITE)


# ============================================================
# Save
# ============================================================
output_path = "/home/felippeb/repos/github/craftura/architecture/THREE-LAYER-ARCHITECTURE.pptx"
prs.save(output_path)
print(f"Slide deck saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
